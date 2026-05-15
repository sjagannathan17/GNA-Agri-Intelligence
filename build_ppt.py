"""
GNA Agri-Intelligence — PowerPoint builder (visual edition)
Z-order rule: background → photos → bars/cards → text → top_bar last
Run from project root:  python3 build_ppt.py
"""
import json, os, io
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image, ImageDraw

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand ──────────────────────────────────────────────────────────────────────
C = dict(
    dark="#2d5016", green="#86AF49", lime="#c8e48a",
    light="#eef4e2", border="#dde8cc", white="#ffffff",
    bg="#f2f5ec", text="#1a1a1a", mid="#333333",
    grey="#888888", red="#f87171", amber="#fbbf24", blue="#60a5fa",
)

def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def mpl(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2],16)/255 for i in (0,2,4))

W, H   = Inches(13.33), Inches(7.5)
M      = Inches(0.52)
BAR_H  = Inches(0.68)

# ── Data ───────────────────────────────────────────────────────────────────────
with open("dashboard/report.json") as f:
    D = json.load(f)
ss, zb, ie  = D["season_summary"], D["zone_breakdown"], D["input_effectiveness"]
pa, ft, pf  = D["priority_actions"], D["farmer_tiers"], D["procurement_funnel"]
fin, sc, rain = D["financial"], D["season_comparison"]["metrics"], D["rainfall"]

IMG = "dashboard/img"
P1 = f"{IMG}/AWright_Zambia_00941-1024x683.jpg"
P2 = f"{IMG}/0D4A7760-1024x683.jpg"
P3 = f"{IMG}/0D4A7934-1024x683.jpg"

def fmt(n):
    if n >= 1e6: return f"{n/1e6:.1f}M"
    if n >= 1e3: return f"{n/1e3:.1f}K"
    return str(int(n))

# ╔═══════════════════════════════════════════════════════════════════════════════
#  HELPERS
# ╚═══════════════════════════════════════════════════════════════════════════════
def rect(sl, x, y, w, h, fill=None, line=None, lw=Pt(0.6), round_=0):
    shtype = 5 if round_ else 1
    sh = sl.shapes.add_shape(shtype, int(x), int(y), int(w), int(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    else:    sh.fill.background()
    if line: sh.line.color.rgb = rgb(line); sh.line.width = lw
    else:    sh.line.fill.background()
    if round_ and shtype == 5:
        sh.adjustments[0] = round_
    return sh

def txt(sl, text, x, y, w, h,
        size=10, bold=False, italic=False,
        color="#1a1a1a", align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = rgb(color); r.font.name = "Calibri"

def lbl(sl, text, x, y, w, color=None, size=7.5):
    tb = sl.shapes.add_textbox(int(x), int(y), int(w), Inches(0.22))
    tf = tb.text_frame; p = tf.paragraphs[0]
    r  = p.add_run(); r.text = text.upper()
    r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = rgb(color or C["green"]); r.font.name = "Calibri"

def embed_image(sl, pil_img, x, y, w, h):
    buf = io.BytesIO()
    pil_img.convert("RGB").save(buf, "JPEG", quality=92)
    buf.seek(0)
    sl.shapes.add_picture(buf, int(x), int(y), int(w), int(h))

def embed_fig(sl, fig, x, y, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0); plt.close(fig)
    sl.shapes.add_picture(buf, int(x), int(y), int(w), int(h))

def load_photo(path, out_px, dark=0.0, crop=None):
    img = Image.open(path).convert("RGB")
    if crop: img = img.crop(crop)
    img = img.resize(out_px, Image.LANCZOS)
    if dark > 0:
        ov = Image.new("RGB", out_px, (int(10*dark), int(30*dark), int(5*dark)))
        img = Image.blend(img, ov, dark)
    return img

def px(inches): return int(inches / Inches(1) * 150)

def hbar(sl, x, y, w, h, pct, fill=None):
    rect(sl, x, y, w, h, fill=C["light"])
    if pct > 0:
        rect(sl, x, y, max(w*min(pct,1), Pt(2)), h, fill=fill or C["green"])

# Big KPI card — label top, huge number, one-line context
def kpi(sl, x, y, w, h, label_, val, sub, vcol=None, bg=C["white"]):
    rect(sl, x, y, w, h, fill=bg, line=C["border"], round_=0.09)
    rect(sl, x, y, w, Inches(0.05), fill=vcol or C["green"])
    txt(sl, label_.upper(), x+Inches(0.18), y+Inches(0.14),
        w-Inches(0.36), Inches(0.22), size=7.5, bold=True,
        color=C["green"] if bg==C["white"] else C["lime"])
    txt(sl, val, x+Inches(0.18), y+Inches(0.38),
        w-Inches(0.36), Inches(0.72), size=34, bold=True,
        color=vcol or C["text"])
    txt(sl, sub, x+Inches(0.18), y+Inches(1.12),
        w-Inches(0.36), Inches(0.22), size=8.5, color=C["grey"])

# ── Top bar (always last) ──────────────────────────────────────────────────────
def top_bar(sl, title, sub=""):
    rect(sl, 0, 0, W, BAR_H, fill=C["dark"])
    txt(sl, title, M, Inches(0.1), Inches(9.5), Inches(0.4),
        size=16, bold=True, color=C["white"])
    if sub:
        txt(sl, sub, M, Inches(0.48), Inches(10), Inches(0.2),
            size=8.5, color=C["lime"])
    rect(sl, W-Inches(1.62), Inches(0.11), Inches(1.52), Inches(0.46),
         fill=C["green"], round_=0.15)
    txt(sl, "GNA  ·  2025/26", W-Inches(1.60), Inches(0.17),
        Inches(1.52), Inches(0.34), size=8, bold=True,
        color=C["white"], align=PP_ALIGN.CENTER)

# ╔═══════════════════════════════════════════════════════════════════════════════
#  MATPLOTLIB CHARTS  (cleaner, more padding, bigger fonts)
# ╚═══════════════════════════════════════════════════════════════════════════════
BG = mpl(C["bg"])

def ax_style(ax):
    ax.set_facecolor(BG)
    ax.spines[:].set_visible(False)
    ax.tick_params(colors=mpl(C["grey"]), labelsize=9, length=0)

def chart_zone_risk():
    zones = sorted(zb, key=lambda z: -z["high_risk_count"])
    names = [f"Zone {z['zone']}" for z in zones]
    vals  = [z["high_risk_count"] for z in zones]
    cols  = [mpl(C["red"]) if z["avg_risk"]>0.45 else
             (mpl(C["amber"]) if z["avg_risk"]>0.35 else mpl(C["green"])) for z in zones]
    fig, ax = plt.subplots(figsize=(6.5, 3.2), facecolor=BG)
    ax_style(ax)
    bars = ax.barh(names, vals, color=cols, height=0.6, zorder=3)
    ax.invert_yaxis()
    ax.grid(axis="x", color=mpl(C["border"]), lw=0.6, zorder=0)
    for b, v in zip(bars, vals):
        ax.text(v+60, b.get_y()+b.get_height()/2,
                f"{v:,}", va="center", fontsize=10,
                color=mpl(C["mid"]), fontweight="bold")
    ax.set_xlim(0, max(vals)*1.25)
    ax.tick_params(axis="y", labelsize=10, colors=mpl(C["dark"]))
    ax.set_xlabel("High-risk farmers", fontsize=9, color=mpl(C["grey"]))
    fig.tight_layout(pad=0.8)
    return fig

def chart_inputs_big():
    keys   = ["inoculant","fertilizer","fungicide","seed_guard"]
    labels = ["Inoculant","Fertilizer","Fungicide","Seed Guard"]
    lifts  = [ie[k]["yield_lift_kg_ha"] for k in keys]
    adopt  = [ie[k]["adoption_pct"]*100 for k in keys]
    cols   = [mpl(C["green"]), mpl(C["blue"]), mpl(C["amber"]), mpl("#a78bfa")]
    fig, ax1 = plt.subplots(figsize=(7.5, 3.4), facecolor=BG)
    ax_style(ax1)
    x = np.arange(len(labels))
    bars = ax1.bar(x, lifts, color=cols, width=0.55, zorder=3)
    ax1.set_ylabel("Extra yield (kg/ha)", fontsize=9, color=mpl(C["dark"]))
    ax1.tick_params(axis="y", labelcolor=mpl(C["dark"]))
    ax2 = ax1.twinx(); ax2.set_facecolor(BG)
    ax2.spines[:].set_visible(False)
    ax2.plot(x, adopt, "o--", color=mpl(C["amber"]), lw=2.5, ms=8, zorder=5)
    ax2.set_ylabel("Adoption %", fontsize=9, color=mpl(C["amber"]))
    ax2.tick_params(axis="y", labelcolor=mpl(C["amber"]), labelsize=9, length=0)
    ax1.set_xticks(x); ax1.set_xticklabels(labels, fontsize=10)
    ax1.grid(axis="y", color=mpl(C["border"]), lw=0.5, zorder=0)
    for b, v in zip(bars, lifts):
        ax1.text(b.get_x()+b.get_width()/2, v+6, f"+{v} kg",
                 ha="center", va="bottom", fontsize=10,
                 color=mpl(C["dark"]), fontweight="bold")
    fig.tight_layout(pad=0.8)
    return fig

def chart_rainfall_big():
    months = ["Nov","Dec","Jan","Feb","Mar","Apr"]
    iia = next(r for r in rain if r["zone"]=="IIa")
    iib = next(r for r in rain if r["zone"]=="IIb")
    x, w2 = np.arange(len(months)), 0.28
    fig, ax = plt.subplots(figsize=(8.0, 3.6), facecolor=BG)
    ax_style(ax)
    ax.bar(x-w2, [iia["historical_monthly_mm"][m] for m in months], w2,
           color=mpl(C["border"]), label="Historical avg", zorder=3)
    ax.bar(x, [iia["monthly_mm"][m] for m in months], w2,
           color=mpl(C["red"]), label="Zone IIa", zorder=3, alpha=0.9)
    ax.bar(x+w2, [iib["monthly_mm"][m] for m in months], w2,
           color=mpl(C["green"]), label="Zone IIb", zorder=3, alpha=0.9)
    ax.set_xticks(x); ax.set_xticklabels(months, fontsize=10)
    ax.set_ylabel("mm", fontsize=9, color=mpl(C["grey"]))
    ax.grid(axis="y", color=mpl(C["border"]), lw=0.5, zorder=0)
    ax.legend(fontsize=9, framealpha=0, labelcolor=mpl(C["mid"]))
    fig.tight_layout(pad=0.8)
    return fig

def chart_risk_trend():
    wt  = D["weekly_risk_trend"]
    hi  = [w["high_risk"] for w in wt]
    wks = [w["week"] for w in wt]
    fig, ax = plt.subplots(figsize=(6.0, 3.0), facecolor=BG)
    ax_style(ax)
    ax.fill_between(range(len(wks)), hi, alpha=0.2, color=mpl(C["red"]))
    ax.plot(range(len(wks)), hi, "o-", color=mpl(C["red"]), lw=2.5, ms=6, zorder=5)
    ax.axhline(5500, color=mpl(C["amber"]), lw=1.4, ls="--")
    ax.text(len(wks)-1, 5580, "watch threshold",
            ha="right", fontsize=9, color=mpl(C["amber"]))
    ax.set_xticks(range(len(wks))); ax.set_xticklabels(wks, fontsize=9)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{int(v):,}"))
    ax.grid(axis="y", color=mpl(C["border"]), lw=0.5, zorder=0)
    fig.tight_layout(pad=0.8)
    return fig

def chart_funnel():
    stages = [s["stage"] for s in pf]
    pcts   = [s["pct"]   for s in pf]
    counts = [s["count"] for s in pf]
    colors = [mpl(C["green"]) if p>0.7 else
              (mpl(C["amber"]) if p>0.6 else mpl(C["red"])) for p in pcts]
    fig, ax = plt.subplots(figsize=(5.5, 3.2), facecolor=BG)
    ax_style(ax)
    y = np.arange(len(stages))
    ax.barh(y, pcts, color=colors, height=0.55, zorder=3)
    ax.set_xlim(0, 1.35)
    ax.set_yticks(y); ax.set_yticklabels(stages, fontsize=9)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0%}"))
    for i,(p,c) in enumerate(zip(pcts,counts)):
        ax.text(p+0.02, i, f"{p:.0%}  ({c:,})",
                va="center", fontsize=9.5, color=mpl(C["mid"]), fontweight="bold")
    ax.invert_yaxis()
    ax.grid(axis="x", color=mpl(C["border"]), lw=0.5, zorder=0)
    fig.tight_layout(pad=0.8)
    return fig

# ╔═══════════════════════════════════════════════════════════════════════════════
#  BUILD
# ╚═══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════════
# S1  TITLE  — full-bleed composite (unchanged — already very visual)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)

SLIDE_PX = (px(W), px(H))
bg = Image.new("RGB", SLIDE_PX, tuple(int(x*255) for x in mpl(C["dark"])))
right_photo = load_photo(P1, (SLIDE_PX[0]//2, SLIDE_PX[1]),
                         dark=0.30, crop=(80, 0, 944, 683))
bg.paste(right_photo, (SLIDE_PX[0]//2, 0))
fade_w = 200
for xi in range(fade_w):
    alpha = 1 - xi / fade_w
    for yi in range(SLIDE_PX[1]):
        ox  = SLIDE_PX[0]//2 + xi
        pr  = right_photo.getpixel((xi, yi))
        dk  = tuple(int(x*255) for x in mpl(C["dark"]))
        blended = tuple(int(dk[c]*alpha + pr[c]*(1-alpha)) for c in range(3))
        bg.putpixel((ox, yi), blended)
embed_image(sl, bg, 0, 0, W, H)

rect(sl, 0, 0, W, Inches(0.055), fill=C["green"])
rect(sl, 0, H-Inches(1.18), W*0.50, Inches(1.18), fill="#1a2e0a")
rect(sl, M, Inches(1.42), Inches(0.85), Inches(0.85), fill=C["green"], round_=0.12)
txt(sl, "GNA", M, Inches(1.44), Inches(0.85), Inches(0.78),
    size=22, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
txt(sl, "GOOD NATURE AGRO", M+Inches(1.0), Inches(1.44),
    Inches(6), Inches(0.3), size=10, bold=True, color=C["green"])
txt(sl, "Zambia · Agri-Intelligence Platform", M+Inches(1.0), Inches(1.76),
    Inches(6), Inches(0.26), size=9, color=C["lime"])
txt(sl, "Season 2025/26",     M, Inches(2.58), Inches(6.5), Inches(0.58),
    size=40, bold=True, color=C["white"])
txt(sl, "Leadership Briefing", M, Inches(3.18), Inches(6.5), Inches(0.58),
    size=40, bold=True, color=C["lime"])
rect(sl, M, Inches(4.5), Inches(1.95), Inches(0.3), fill=C["green"], round_=0.3)
txt(sl, "May 2026  ·  Confidential",
    M+Inches(0.1), Inches(4.54), Inches(1.95), Inches(0.25),
    size=8, bold=True, color=C["white"])

stats = [(fmt(ss["total_farmers"]), "active farmers"),
         (f"{ss['buyback_rate']:.0%}", "buyback rate"),
         (fmt(ss["high_risk_count"]), "high-risk farmers")]
SW = W*0.50/3
for i,(v,l) in enumerate(stats):
    sx = i*SW
    if i: rect(sl, sx, H-Inches(1.08), Pt(0.5), Inches(0.72), fill=C["border"])
    txt(sl, v, sx+Inches(0.18), H-Inches(1.06),
        SW-Inches(0.2), Inches(0.52), size=22, bold=True, color=C["green"])
    txt(sl, l, sx+Inches(0.18), H-Inches(0.52),
        SW-Inches(0.2), Inches(0.24), size=8, color=C["lime"])

txt(sl, "Photo: Good Nature Agro / Zambia",
    W*0.54, H-Inches(0.28), W*0.44, Inches(0.22),
    size=6.5, italic=True, color="#aaaaaa", align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# S2  SEASON AT A GLANCE  — PIL composite bg (photo right) + KPI cards left
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)

# Bake slide bg as PIL composite so photo is guaranteed visible
SPX2 = (px(W), px(H))
bg2  = Image.new("RGB", SPX2, tuple(int(x*255) for x in mpl(C["bg"])))
# Right 43% = farmer photo
PW2  = int(SPX2[0] * 0.43)
PX2  = SPX2[0] - PW2
ph2  = load_photo(P2, (PW2, SPX2[1]), dark=0.10, crop=(80, 0, 944, 683))
bg2.paste(ph2, (PX2, 0))
# Gradient fade at the join (150px)
FW2 = 150
for xi in range(FW2):
    a2 = 1 - xi / FW2
    for yi in range(SPX2[1]):
        pr2 = ph2.getpixel((xi, yi))
        bk2 = tuple(int(x*255) for x in mpl(C["bg"]))
        bl2 = tuple(int(bk2[c]*a2 + pr2[c]*(1-a2)) for c in range(3))
        bg2.putpixel((PX2+xi, yi), bl2)
embed_image(sl, bg2, 0, 0, W, H)   # single embed — baked in, always visible

# 4 tall KPI cards on top of the baked bg
CONTENT_Y = BAR_H + Inches(0.3)
CONTENT_H = H - CONTENT_Y - Inches(0.22)
CARD_W    = (W * 0.55 - 2*M - Inches(0.15)) / 2
KH2 = (CONTENT_H - Inches(0.18)) / 2

kpi_data = [
    ("Active Farmers",    fmt(ss["total_farmers"]),        "enrolled this season",   C["dark"]),
    ("Buyback Rate",      f"{ss['buyback_rate']:.0%}",     "sold back to GNA",       C["green"]),
    ("High-Risk Farmers", fmt(ss["high_risk_count"]),      "need field visit now",   C["red"]),
    ("In-Kind Gap Rate",  f"{ss['in_kind_gap_rate']:.0%}", "cannot cover repayment", C["amber"]),
]
for i,(l,v,s,c) in enumerate(kpi_data):
    cx2 = M + (i%2)*(CARD_W+Inches(0.15))
    cy2 = CONTENT_Y + (i//2)*(KH2+Inches(0.18))
    kpi(sl, cx2, cy2, CARD_W, KH2, l, v, s, vcol=c)

# Photo credit caption over the photo area
txt(sl, "GNA farmer · Zambia",
    W*0.6, H-Inches(0.32), W*0.38, Inches(0.22),
    size=7, italic=True, color="#cccccc", align=PP_ALIGN.RIGHT)

top_bar(sl, "Season at a Glance", "2025/26 · Zambia · 22,597 farmers")

# ══════════════════════════════════════════════════════════════════════════════
# S3  ZONE RISK  — 5 visual zone cards (top) + bar chart (bottom)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=C["bg"])

ZN = {"I":"Gwembe & Luangwa","IIa":"Central Plateau","IIb":"Southern Plateau",
      "III":"Northern Province","IV":"Western Barotse"}

# Zone cards — top half
CARD_Y = BAR_H + Inches(0.28)
CARD_H = Inches(2.9)
CARD_W = (W - 2*M - Inches(0.8)) / 5

for i, z in enumerate(sorted(zb, key=lambda z: z["zone"])):
    cx = M + i*(CARD_W+Inches(0.2))
    scol = C["red"] if z["avg_risk"]>0.45 else (C["amber"] if z["avg_risk"]>0.35 else C["green"])
    sbg  = "#fef2f2" if scol==C["red"] else ("#fffbeb" if scol==C["amber"] else C["light"])

    rect(sl, cx, CARD_Y, CARD_W, CARD_H, fill=sbg, line=C["border"], round_=0.08)
    # Colored top band
    rect(sl, cx, CARD_Y, CARD_W, Inches(0.065), fill=scol)

    # Zone label
    txt(sl, f"ZONE {z['zone']}", cx+Inches(0.14), CARD_Y+Inches(0.14),
        CARD_W-Inches(0.28), Inches(0.24), size=9.5, bold=True, color=scol)
    txt(sl, ZN.get(z["zone"],""), cx+Inches(0.14), CARD_Y+Inches(0.37),
        CARD_W-Inches(0.28), Inches(0.2), size=7.5, color=C["grey"])

    # Big risk score
    txt(sl, f"{z['avg_risk']:.2f}", cx+Inches(0.14), CARD_Y+Inches(0.6),
        CARD_W-Inches(0.28), Inches(0.72), size=42, bold=True, color=scol)
    txt(sl, "risk score", cx+Inches(0.14), CARD_Y+Inches(1.32),
        CARD_W-Inches(0.28), Inches(0.2), size=8, color=C["grey"])

    # Status badge
    status = "URGENT" if z["avg_risk"]>0.45 else ("WATCH" if z["avg_risk"]>0.35 else "OK")
    rect(sl, cx+Inches(0.14), CARD_Y+Inches(1.58),
         CARD_W-Inches(0.28), Inches(0.26), fill=scol, round_=0.3)
    txt(sl, status, cx+Inches(0.14), CARD_Y+Inches(1.58),
        CARD_W-Inches(0.28), Inches(0.26), size=8.5, bold=True,
        color=C["white"], align=PP_ALIGN.CENTER)

    # Farmer count + at-risk
    txt(sl, f"{z['count']:,}",
        cx+Inches(0.14), CARD_Y+Inches(1.96),
        CARD_W-Inches(0.28), Inches(0.3), size=14, bold=True, color=C["text"])
    txt(sl, "farmers", cx+Inches(0.14), CARD_Y+Inches(2.26),
        CARD_W-Inches(0.28), Inches(0.18), size=7.5, color=C["grey"])

    # Inoculant bar
    hbar(sl, cx+Inches(0.14), CARD_Y+Inches(2.5),
         CARD_W-Inches(0.28), Inches(0.1), z["inoculant_pct"], fill=C["green"])
    txt(sl, f"Inoculant {z['inoculant_pct']:.0%}",
        cx+Inches(0.14), CARD_Y+Inches(2.64),
        CARD_W-Inches(0.28), Inches(0.18), size=7, color=C["grey"])

# Bar chart — bottom
BY3 = CARD_Y + CARD_H + Inches(0.32)
lbl(sl, "High-Risk Farmer Count by Zone", M, BY3-Inches(0.24), W-2*M)
embed_fig(sl, chart_zone_risk(), M, BY3, W-2*M, H-BY3-Inches(0.18))

top_bar(sl, "Zone Risk Overview",
        "Risk scores · 22,597 farmers · 5 zones · Season 2025/26")

# ══════════════════════════════════════════════════════════════════════════════
# S4  INPUT EFFECTIVENESS  — 4 big input cards + full-width chart
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=C["bg"])

IY = BAR_H + Inches(0.28)
IW = (W-2*M-Inches(0.45))/4
IH = Inches(2.2)
INPUT_META = {
    "inoculant":  ("Inoculant",  C["green"]),
    "fertilizer": ("Fertilizer", C["blue"]),
    "fungicide":  ("Fungicide",  C["amber"]),
    "seed_guard": ("Seed Guard", "#a78bfa"),
}
for i, key in enumerate(["inoculant","fertilizer","fungicide","seed_guard"]):
    d = ie[key]; nm, col = INPUT_META[key]
    cx = M + i*(IW+Inches(0.15))
    rect(sl, cx, IY, IW, IH, fill=C["white"], line=C["border"], round_=0.09)
    rect(sl, cx, IY, IW, Inches(0.06), fill=col)

    txt(sl, nm.upper(), cx+Inches(0.18), IY+Inches(0.14),
        IW-Inches(0.36), Inches(0.24), size=9, bold=True, color=C["dark"])

    # Giant yield lift number
    txt(sl, f"+{d['yield_lift_kg_ha']}",
        cx+Inches(0.14), IY+Inches(0.42),
        IW-Inches(0.28), Inches(0.82), size=46, bold=True, color=col)
    txt(sl, "kg/ha", cx+Inches(0.18), IY+Inches(1.26),
        IW-Inches(0.36), Inches(0.24), size=10, color=C["grey"])

    # Adoption bar + label
    hbar(sl, cx+Inches(0.18), IY+Inches(1.62),
         IW-Inches(0.36), Inches(0.12), d["adoption_pct"], fill=col)
    txt(sl, f"{d['adoption_pct']:.0%} adoption",
        cx+Inches(0.18), IY+Inches(1.8),
        IW-Inches(0.36), Inches(0.24), size=9.5, bold=True, color=col)

# Full-width chart below
CY4 = IY + IH + Inches(0.34)
lbl(sl, "Yield Lift vs. Adoption Rate per Input", M, CY4-Inches(0.24), W-2*M)
embed_fig(sl, chart_inputs_big(), M, CY4, W-2*M, H-CY4-Inches(0.18))

top_bar(sl, "Input Effectiveness",
        "Measured yield lift across GNA farmers · Season 2025/26")

# ══════════════════════════════════════════════════════════════════════════════
# S5  RAINFALL  — 5 big zone cards + full-width chart
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=C["bg"])

RY = BAR_H + Inches(0.28)
DW = (W-2*M-Inches(0.8))/5
DH = Inches(1.9)

for i, r in enumerate(rain):
    cx   = M + i*(DW+Inches(0.2))
    dcol = C["red"] if r["anomaly_pct"]<-10 else (C["amber"] if r["anomaly_pct"]<-5 else C["green"])
    dbg  = "#fef2f2" if dcol==C["red"] else ("#fffbeb" if dcol==C["amber"] else C["light"])

    rect(sl, cx, RY, DW, DH, fill=dbg, line=C["border"], round_=0.09)
    rect(sl, cx, RY, DW, Inches(0.06), fill=dcol)

    txt(sl, f"ZONE {r['zone']}", cx+Inches(0.14), RY+Inches(0.14),
        DW-Inches(0.28), Inches(0.22), size=9, bold=True, color=dcol)

    # Giant anomaly number
    txt(sl, f"{r['anomaly_pct']:+.0f}%",
        cx+Inches(0.12), RY+Inches(0.4),
        DW-Inches(0.24), Inches(0.78), size=38, bold=True, color=dcol)
    txt(sl, "vs avg",
        cx+Inches(0.14), RY+Inches(1.2),
        DW-Inches(0.28), Inches(0.2), size=8, color=C["grey"])

    # Progress bar (relative to historical)
    hbar(sl, cx+Inches(0.14), RY+Inches(1.52),
         DW-Inches(0.28), Inches(0.1),
         1+r["anomaly_pct"]/100, fill=dcol)
    txt(sl, f"{r['season_total_mm']} mm",
        cx+Inches(0.14), RY+Inches(1.68),
        DW-Inches(0.28), Inches(0.18), size=8.5, bold=True, color=C["text"])

# Full-width chart below
CHTY = RY + DH + Inches(0.32)
lbl(sl, "Monthly Rainfall  ·  Zone IIa vs Zone IIb vs Historical Average",
    M, CHTY-Inches(0.24), W-2*M)
embed_fig(sl, chart_rainfall_big(), M, CHTY, W-2*M, H-CHTY-Inches(0.18))

top_bar(sl, "Rainfall & Zone Conditions",
        "Actual vs. historical average · Nov 2025 – Apr 2026")

# ══════════════════════════════════════════════════════════════════════════════
# S6  FARMER SEGMENTS & FINANCIAL  — visual tiers + financial KPIs + charts
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=C["bg"])

COL_W = (W-2*M-Inches(0.6))/3
LX6  = M
MX6  = M + COL_W + Inches(0.3)
RX6  = M + 2*(COL_W+Inches(0.3))
SY6  = BAR_H + Inches(0.28)

# LEFT: Farmer tier cards (visual, no description text)
lbl(sl, "Farmer Tiers", LX6, SY6-Inches(0.24), COL_W)
TIER_META = {
    "Elite":   (C["green"], C["light"]),
    "Core":    (C["blue"],  "#eff6ff"),
    "Growth":  ("#a78bfa",  "#f5f3ff"),
    "At-Risk": (C["red"],   "#fef2f2"),
}
TH6 = (H - SY6 - Inches(0.22)) / 4 - Inches(0.1)
TY6 = SY6
for t in ft:
    acc, tbg = TIER_META[t["tier"]]
    rect(sl, LX6, TY6, COL_W, TH6, fill=tbg, line=C["border"], round_=0.07)
    rect(sl, LX6, TY6, Inches(0.05), TH6, fill=acc)

    txt(sl, t["tier"].upper(),
        LX6+Inches(0.14), TY6+Inches(0.1),
        COL_W*0.4, Inches(0.26), size=11, bold=True, color=acc)

    # Volume % — big
    txt(sl, f"{t['volume_pct']:.0%}",
        LX6+COL_W*0.46, TY6+Inches(0.06),
        COL_W*0.28, TH6-Inches(0.12), size=28, bold=True, color=acc)

    txt(sl, f"{t['farmers']:,} farmers",
        LX6+Inches(0.14), TY6+TH6-Inches(0.28),
        COL_W*0.55, Inches(0.2), size=8.5, color=C["grey"])
    txt(sl, "of volume",
        LX6+COL_W*0.46, TY6+TH6-Inches(0.28),
        COL_W*0.5, Inches(0.2), size=8, color=C["grey"])
    TY6 += TH6 + Inches(0.1)

# MIDDLE: Procurement funnel chart
lbl(sl, "Procurement Funnel", MX6, SY6-Inches(0.24), COL_W)
embed_fig(sl, chart_funnel(), MX6, SY6, COL_W, H-SY6-Inches(0.22))

# RIGHT: 3 big financial KPIs + risk trend chart
lbl(sl, "Financial Exposure", RX6, SY6-Inches(0.24), COL_W)
fin3 = [
    ("Total Loan Book",  f"${fin['total_loan_book_usd']/1e6:.1f}M",  C["dark"]),
    ("At-Risk Exposure", f"${fin['at_risk_loan_usd']/1e6:.1f}M",     C["red"]),
    ("Exposure Gap",     f"${fin['gap_usd']/1e6:.1f}M",              C["red"]),
]
FH6 = Inches(1.1)
FY6 = SY6
for l6,v6,c6 in fin3:
    rect(sl, RX6, FY6, COL_W, FH6, fill=C["white"], line=C["border"], round_=0.07)
    rect(sl, RX6, FY6, COL_W, Inches(0.05), fill=c6)
    txt(sl, l6.upper(), RX6+Inches(0.14), FY6+Inches(0.12),
        COL_W-Inches(0.28), Inches(0.2), size=7.5, bold=True, color=C["grey"])
    txt(sl, v6, RX6+Inches(0.14), FY6+Inches(0.34),
        COL_W-Inches(0.28), Inches(0.52), size=28, bold=True, color=c6)
    FY6 += FH6 + Inches(0.1)

RTY6 = FY6 + Inches(0.22)
lbl(sl, "Weekly Risk Trend", RX6, RTY6-Inches(0.24), COL_W)
embed_fig(sl, chart_risk_trend(), RX6, RTY6, COL_W, H-RTY6-Inches(0.22))

top_bar(sl, "Farmer Segments & Financial Exposure", "Season 2025/26")

# ══════════════════════════════════════════════════════════════════════════════
# S7  PRIORITY ACTIONS  — headline + impact, minimal text
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=C["bg"])

URGENCY = {
    "critical": (C["red"],   "CRITICAL"),
    "high":     (C["amber"], "HIGH"),
    "medium":   (C["blue"],  "MEDIUM"),
}
AY7 = BAR_H + Inches(0.26)
AH7 = (H - AY7 - Inches(0.22)) / len(pa) - Inches(0.1)

for action in pa:
    col7, lbl7 = URGENCY[action["urgency"]]

    rect(sl, M, AY7, W-2*M, AH7, fill=C["white"], line=C["border"], round_=0.07)
    # Left accent bar
    rect(sl, M, AY7, Inches(0.06), AH7, fill=col7)

    # Priority circle
    cx7 = M + Inches(0.22)
    cy7 = AY7 + (AH7 - Inches(0.56))/2
    rect(sl, cx7, cy7, Inches(0.56), Inches(0.56), fill=col7, round_=0.5)
    txt(sl, str(action["priority"]), cx7, cy7,
        Inches(0.56), Inches(0.56), size=18, bold=True,
        color=C["white"], align=PP_ALIGN.CENTER)

    # Urgency badge
    BX7 = M + Inches(0.94)
    BY7 = AY7 + Inches(0.14)
    rect(sl, BX7, BY7, Inches(0.82), Inches(0.28), fill=col7, round_=0.4)
    txt(sl, lbl7, BX7, BY7, Inches(0.82), Inches(0.28),
        size=8, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    # Action headline — truncated to one line (no paragraph wrapping)
    headline = action["action"].split(".")[0]  # first sentence only
    if len(headline) > 80: headline = headline[:77] + "…"
    txt(sl, headline, BX7, AY7 + Inches(0.5),
        W - 2*M - Inches(3.4), AH7 - Inches(0.62),
        size=10.5, bold=False, color=C["mid"], wrap=False)

    # Impact pill — bigger, prominent
    IPW7 = Inches(2.8)
    IPH7 = Inches(0.56)
    IPX7 = W - M - IPW7
    IPY7 = AY7 + (AH7 - IPH7)/2
    rect(sl, IPX7, IPY7, IPW7, IPH7, fill=C["light"], line=C["border"], round_=0.25)
    txt(sl, action["impact"], IPX7, IPY7, IPW7, IPH7,
        size=9.5, bold=True, color=C["dark"], align=PP_ALIGN.CENTER)

    AY7 += AH7 + Inches(0.1)

top_bar(sl, "Priority Actions",
        "AI-ranked by procurement impact · Week ending May 3, 2026")

# ══════════════════════════════════════════════════════════════════════════════
# S8  WHATSAPP ADVISOR  — PIL composite bg (farmer photo left strip) + phone
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)

# Bake bg: full photo with strong dark overlay so text is legible
SPX8 = (px(W), px(H))
bg8  = load_photo(P1, SPX8, dark=0.72, crop=(0, 0, 1024, 683))
# Left 55%: slightly darker for content area
lw8 = int(SPX8[0] * 0.56)
for xi in range(lw8):
    extra = max(0, 1 - xi / (lw8 * 0.7)) * 0.35
    for yi in range(SPX8[1]):
        pr8 = bg8.getpixel((xi, yi))
        bg8.putpixel((xi, yi), tuple(int(pr8[c]*(1-extra)) for c in range(3)))
embed_image(sl, bg8, 0, 0, W, H)

HW8  = (W-2*M-Inches(0.36))/2
LX8  = M
RX8  = M + HW8 + Inches(0.36)
WY8  = BAR_H + Inches(0.28)

# LEFT: How it works — big numbered circles, headline only (white text on dark photo bg)
lbl(sl, "How It Works", LX8, WY8-Inches(0.24), HW8, color=C["lime"])
steps8 = [
    ("1", C["green"],  "Farmer sends a WhatsApp message"),
    ("2", C["blue"],   "AI reads live risk profile"),
    ("3", C["amber"],  "Claude generates personalised reply"),
    ("4", C["red"],    "Field agent alerted if high-risk"),
]
STEP_H = Inches(0.88)
for i,(num,col8,hd) in enumerate(steps8):
    sy8 = WY8 + i*(STEP_H + Inches(0.12))
    rect(sl, LX8, sy8, STEP_H, STEP_H, fill=col8, round_=0.5)
    txt(sl, num, LX8, sy8, STEP_H, STEP_H,
        size=24, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    txt(sl, hd, LX8+STEP_H+Inches(0.2), sy8+Inches(0.22),
        HW8-STEP_H-Inches(0.3), Inches(0.44),
        size=12, bold=True, color=C["white"])

# Engagement KPIs — dark cards so numbers pop against photo bg
EY8 = WY8 + len(steps8)*(STEP_H+Inches(0.12)) + Inches(0.36)
lbl(sl, "Season Engagement", LX8, EY8-Inches(0.24), HW8, color=C["lime"])
EW8 = (HW8-Inches(0.3))/3
for i,(l,v,s) in enumerate([
    ("Response Rate", "63%",   "nudge engagement"),
    ("Msgs / Farmer", "7.3",   "avg over season"),
    ("Escalations",   "1,842", "to field agents"),
]):
    ex8 = LX8 + i*(EW8+Inches(0.15))
    rect(sl, ex8, EY8, EW8, Inches(1.2), fill="#1a3008", line=C["green"], lw=Pt(1), round_=0.09)
    txt(sl, l.upper(), ex8+Inches(0.12), EY8+Inches(0.12),
        EW8-Inches(0.24), Inches(0.2), size=7, bold=True, color=C["lime"])
    txt(sl, v, ex8+Inches(0.12), EY8+Inches(0.34),
        EW8-Inches(0.24), Inches(0.5), size=26, bold=True, color=C["green"])
    txt(sl, s, ex8+Inches(0.12), EY8+Inches(0.9),
        EW8-Inches(0.24), Inches(0.2), size=8, color=C["lime"])

# RIGHT: Phone mockup
lbl(sl, "Live Conversation  ·  Agnes Banda, Zone IIa", RX8, WY8-Inches(0.24), HW8, color=C["lime"])
PH8 = H - WY8 - Inches(0.2)
rect(sl, RX8, WY8, HW8, PH8, fill="#111b21", line="#2a3942", lw=Pt(1))
rect(sl, RX8, WY8, HW8, Inches(0.56), fill="#202c33")
avatar = load_photo(P1, (60, 60), dark=0.1, crop=(280,80,680,560))
embed_image(sl, avatar, RX8+Inches(0.12), WY8+Inches(0.07),
            Inches(0.42), Inches(0.42))
txt(sl, "Agnes Banda",
    RX8+Inches(0.64), WY8+Inches(0.06),
    HW8-Inches(0.78), Inches(0.26), size=10, bold=True, color="#e9edef")
txt(sl, "online · Zone IIa · Season 1",
    RX8+Inches(0.64), WY8+Inches(0.32),
    HW8-Inches(0.78), Inches(0.18), size=8, color="#00a884")

msgs8 = [
    ("in",  "My soybean leaves are turning yellow. What should I do?"),
    ("out", "Hi Agnes! 🌱 Yellow leaves = nitrogen deficiency.\n"
            "1️⃣ Apply inoculant today\n"
            "2️⃣ Add top-dress if available\nReply 1 or 2 ✅"),
    ("in",  "I have inoculant at home but haven't applied it"),
    ("out", "Apply today — dissolve in water, drench roots.\n"
            "Agent Joseph will call tomorrow to confirm. ✅"),
]
MY8 = WY8 + Inches(0.62)
for sender8, text8 in msgs8:
    is_in8 = sender8 == "in"
    lines8 = max(1, len(text8)//40)
    mh8    = Inches(0.26*(lines8+1)+0.1)
    MX8    = RX8+Inches(0.12) if is_in8 else RX8+Inches(0.44)
    MW8    = HW8*0.84
    rect(sl, MX8, MY8, MW8, mh8, fill="#202c33" if is_in8 else "#005c4b")
    txt(sl, text8, MX8+Inches(0.1), MY8+Inches(0.07),
        MW8-Inches(0.2), mh8-Inches(0.1),
        size=8.5, color="#e9edef", wrap=True)
    MY8 += mh8 + Inches(0.1)

top_bar(sl, "WhatsApp Farmer Advisor",
        "AI-powered · real-time · personalised")

# ══════════════════════════════════════════════════════════════════════════════
# S9  CLOSING — full-bleed photo
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)

SLIDE_PX9 = (px(W), px(H))
base9 = load_photo(P3, SLIDE_PX9, dark=0.55)
lw9 = SLIDE_PX9[0]//2
for xi in range(lw9+100):
    alpha = max(0, 1 - xi/(lw9*0.85))
    for yi in range(SLIDE_PX9[1]):
        pr9 = base9.getpixel((xi,yi))
        blended9 = tuple(int(pr9[c]*(1-alpha*0.65)) for c in range(3))
        base9.putpixel((xi,yi), blended9)
embed_image(sl, base9, 0, 0, W, H)

rect(sl, 0, 0, W, Inches(0.055), fill=C["green"])
rect(sl, 0, H-Inches(1.1), W, Inches(1.1), fill="#122407")
rect(sl, M, Inches(1.42), Inches(0.82), Inches(0.82), fill=C["green"], round_=0.1)
txt(sl, "GNA", M, Inches(1.44), Inches(0.82), Inches(0.76),
    size=20, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
txt(sl, "GOOD NATURE AGRO", M+Inches(0.96), Inches(1.47),
    Inches(6), Inches(0.3), size=10, bold=True, color=C["green"])
txt(sl, "Agri-Intelligence Platform",
    M, Inches(2.34), Inches(6.5), Inches(0.6),
    size=30, bold=True, color=C["white"])
txt(sl, "Season 2025/26  ·  Zambia",
    M, Inches(3.04), Inches(6), Inches(0.35), size=12, color=C["lime"])
txt(sl, "Helping Zambian farmers reach the middle class",
    M, Inches(3.52), Inches(6.5), Inches(0.35),
    size=11, italic=True, color=C["lime"])
txt(sl, "goodnatureagro.com  ·  Chipata, Zambia  ·  Photos: Good Nature Agro",
    M, H-Inches(0.66), W-2*M, Inches(0.26), size=8, color=C["grey"])
txt(sl, "Built with GNA Agri-Intelligence · Powered by Claude AI",
    M, H-Inches(0.36), W-2*M, Inches(0.26), size=8, bold=True, color=C["green"])

# ── Save ───────────────────────────────────────────────────────────────────────
out = "dashboard/GNA_AgriIntelligence_Season2025-26.pptx"
prs.save(out)
print(f"✓  {out}  ({len(prs.slides)} slides, {os.path.getsize(out)//1024} KB)")
